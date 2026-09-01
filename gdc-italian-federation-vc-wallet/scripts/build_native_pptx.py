#!/usr/bin/env python3
"""Native, editable PPTX from deck.md + it-wallet-dtd style (not Marp screenshots)."""

from __future__ import annotations

import re
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "deck.md"
OUT = ROOT / "deck.pptx"

BRAND = RGBColor(0x00, 0x66, 0xCC)
BRAND_DARK = RGBColor(0x00, 0x4C, 0x99)
BRAND_SOFT = RGBColor(0xE8, 0xF3, 0xFC)
BRAND_LINE = RGBColor(0xD6, 0xE6, 0xF5)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x6F, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Titillium Web"
FONT_FALLBACK = "Liberation Sans"
MONO = "Consolas"

SW = Inches(13.333333)
SH = Inches(7.5)
ML = Inches(0.58)
MR = Inches(0.50)
CONTENT_W = SW - ML - MR


def _font_name() -> str:
    return FONT


def rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def set_run(run, *, size, color, bold=False, italic=False, mono=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = MONO if mono else _font_name()
    rPr = run._r.get_or_add_rPr()
    # latin + ea so PowerPoint does not swap to Calibri
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", MONO if mono else _font_name())


INLINE_RE = re.compile(
    r"(\*\*(.+?)\*\*|`([^`]+)`|\[([^\]]+)\]\(([^)]+)\)"
    r"|(?<!\*)\*(?!\*)([^*]+?)(?<!\*)\*(?!\*)"
    r"|_([^_]+?)_)"
)


def add_inline(paragraph, text: str, *, size: float, color=TEXT, bold=False):
    pos = 0
    for m in INLINE_RE.finditer(text):
        if m.start() > pos:
            run = paragraph.add_run()
            run.text = text[pos : m.start()]
            set_run(run, size=size, color=color, bold=bold)
        if m.group(2) is not None:
            run = paragraph.add_run()
            run.text = m.group(2)
            set_run(run, size=size, color=color, bold=True)
        elif m.group(3) is not None:
            run = paragraph.add_run()
            run.text = m.group(3)
            set_run(run, size=size, color=BRAND_DARK, bold=False, mono=True)
        elif m.group(4) is not None:
            run = paragraph.add_run()
            run.text = m.group(4)
            set_run(run, size=size, color=BRAND, bold=True)
            run.hyperlink.address = m.group(5)
        elif m.group(6) is not None or m.group(7) is not None:
            run = paragraph.add_run()
            run.text = m.group(6) or m.group(7)
            set_run(run, size=size, color=color, italic=True)
        pos = m.end()
    if pos < len(text):
        run = paragraph.add_run()
        run.text = text[pos:]
        set_run(run, size=size, color=color, bold=bold)


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def clear_first_p(tf):
    p = tf.paragraphs[0]
    p.clear()
    return p


def set_slide_fill(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_picture_fit(slide, path: Path, left, top, width=None, height=None):
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def content_logo(slide):
    logo = ROOT / "images" / "dtd-logo-blue.png"
    add_picture_fit(slide, logo, SW - Inches(3.70), Inches(0.26), width=Inches(3.20))


def content_footer(slide, footer: str):
    box, tf = add_textbox(slide, ML, SH - Inches(0.38), CONTENT_W, Inches(0.30))
    p = clear_first_p(tf)
    run = p.add_run()
    run.text = footer
    set_run(run, size=9, color=MUTED)


def content_title(slide, title: str, *, size=26):
    box, tf = add_textbox(slide, ML, Inches(0.68), CONTENT_W - Inches(3.3), Inches(0.55))
    p = clear_first_p(tf)
    run = p.add_run()
    run.text = title
    set_run(run, size=size, color=BRAND, bold=True)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ML, Inches(1.22), CONTENT_W, Emu(12700)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND
    line.line.fill.background()
    return Inches(1.38)


def chrome(slide, title: str, footer: str, *, title_size=26):
    set_slide_fill(slide, WHITE)
    content_logo(slide)
    content_footer(slide, footer)
    return content_title(slide, title, size=title_size)


def set_bullet(paragraph, *, color=BRAND):
    pPr = paragraph._p.get_or_add_pPr()
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", "Arial")
    buClr = etree.SubElement(pPr, qn("a:buClr"))
    srgb = etree.SubElement(buClr, qn("a:srgbClr"))
    srgb.set("val", rgb_hex(color))
    etree.SubElement(pPr, qn("a:buChar")).set("char", "•")


def add_bullets(tf, items, *, size, first=True, space_after=6):
    for i, item in enumerate(items):
        p = clear_first_p(tf) if (first and i == 0) else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(space_after)
        p.line_spacing = 1.18
        set_bullet(p)
        add_inline(p, item, size=size)


def style_table(table, *, header=True, font=12, col0_bold=True):
    tbl = table._tbl
    tblPr = tbl.tblPr if tbl.tblPr is not None else etree.SubElement(tbl, qn("a:tblPr"))
    # drop theme banding so our fills win
    for child in list(tblPr):
        if child.tag.endswith("tblStyle"):
            tblPr.remove(child)
    for r_idx, row in enumerate(table.rows):
        for c_idx, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            for p in tf.paragraphs:
                for run in p.runs:
                    is_header = header and r_idx == 0
                    set_run(
                        run,
                        size=font,
                        color=WHITE if is_header else TEXT,
                        bold=is_header or (col0_bold and c_idx == 0),
                    )
            fill = cell.fill
            fill.solid()
            if header and r_idx == 0:
                fill.fore_color.rgb = BRAND
            elif r_idx % 2 == 0:
                fill.fore_color.rgb = BRAND_SOFT
            else:
                fill.fore_color.rgb = WHITE
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                el = etree.SubElement(tcPr, qn(side))
                el.set("w", "6350")
                etree.SubElement(el, qn("a:solidFill")).append(
                    etree.Element(qn("a:srgbClr"), val="D6E6F5")
                )


def fill_cell(cell, text: str, *, size, header=False, bold=False):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.clear()
    add_inline(p, text, size=size, color=WHITE if header else TEXT, bold=bold or header)


def parse_front_matter(md: str) -> tuple[dict, str]:
    if not md.startswith("---"):
        return {}, md
    parts = md.split("\n---\n", 1)
    if len(parts) != 2:
        return {}, md
    meta = {}
    for line in parts[0].splitlines()[1:]:
        if ":" in line:
            k, v = line.split(":", 1)
            meta[k.strip()] = v.strip().strip("'\"")
    return meta, parts[1]


def split_slides(body: str) -> list[dict]:
    chunks = re.split(r"\n---\n", body)
    slides = []
    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk:
            continue
        classes = []
        m = re.search(r"<!--\s*_class:\s*(.*?)\s*-->", chunk)
        if m:
            classes = m.group(1).split()
            chunk = (chunk[: m.start()] + chunk[m.end() :]).strip()
        title = ""
        hm = re.search(r"^#{1,2}\s+(.+)$", chunk, re.M)
        if hm:
            title = re.sub(r"\*(.+?)\*", r"\1", hm.group(1).strip())
            chunk = (chunk[: hm.start()] + chunk[hm.end() :]).strip()
        slides.append({"classes": classes, "title": title, "body": chunk})
    return slides


def strip_md_html(text: str) -> str:
    return re.sub(r"</?[^>]+>", "", text).strip()


def parse_lists(body: str) -> list[str]:
    items = []
    for line in body.splitlines():
        m = re.match(r"^\s*(?:[-*]|\d+\.)\s+(.+)$", line)
        if m:
            items.append(m.group(1).strip())
    return items


def parse_table(body: str) -> tuple[list[str], list[list[str]]] | None:
    rows = []
    for line in body.splitlines():
        line = line.strip()
        if line.startswith("|") and not re.match(r"^\|[\s\-:|]+\|$", line):
            cells = [c.strip() for c in line.strip("|").split("|")]
            rows.append(cells)
    if len(rows) < 2:
        return None
    return rows[0], rows[1:]


def parse_stats(body: str) -> list[tuple[str, str]]:
    nums = re.findall(r'class="stat-n">([^<]+)', body)
    labs = re.findall(r'class="stat-l">([^<]+)', body)
    return list(zip(nums, labs))


def parse_images(body: str) -> list[tuple[str, str | None]]:
    found = []
    for m in re.finditer(r"!\[([^\]]*)\]\(([^)]+)\)", body):
        found.append((m.group(2), None))
    for m in re.finditer(
        r'<a href="([^"]+)"[^>]*>\s*<img[^>]+src="([^"]+)"', body, re.S
    ):
        found.append((m.group(2), m.group(1)))
    for m in re.finditer(r'<img[^>]+src="([^"]+)"', body):
        if not any(m.group(1) == p for p, _ in found):
            found.append((m.group(1), None))
    return found


def paragraphs(body: str) -> list[str]:
    text = re.sub(r"<div[\s\S]*?</div>", "\n", body)
    text = re.sub(r"<img[^>]*>", "", text)
    text = re.sub(r"<a href=\"[^\"]+\">\s*</a>", "", text)
    blocks = []
    for raw in re.split(r"\n\s*\n", text):
        raw = raw.strip()
        if not raw:
            continue
        if raw.startswith("|") or raw.startswith("- ") or re.match(r"^\d+\. ", raw):
            continue
        if raw.startswith("![") or raw.startswith("<"):
            continue
        line = " ".join(ln.strip() for ln in raw.splitlines() if not ln.startswith("|"))
        line = strip_md_html(line)
        if line:
            blocks.append(line)
    return blocks


def mermaid_png(mmd: Path, dest: Path) -> Path:
    dest.parent.mkdir(parents=True, exist_ok=True)
    puppeteer_cfg = ROOT / "scripts" / "puppeteer-no-sandbox.json"
    cmd = [
        "npx",
        "--yes",
        "@mermaid-js/mermaid-cli@latest",
        "-p",
        str(puppeteer_cfg),
        "-i",
        str(mmd),
        "-o",
        str(dest),
        "-b",
        "transparent",
        "-s",
        "2",
    ]
    subprocess.run(cmd, check=True, cwd=str(mmd.parent))
    return dest


def resolve_asset(src: str, tmp: Path) -> Path:
    path = (ROOT / src).resolve()
    if path.suffix.lower() == ".svg":
        mmd = path.with_suffix(".mmd")
        png = tmp / (path.stem + ".png")
        if mmd.exists():
            return mermaid_png(mmd, png)
        subprocess.run(
            [
                "inkscape",
                str(path),
                "--export-type=png",
                f"--export-filename={png}",
                "--export-width=1800",
                "--export-background-opacity=0",
            ],
            check=True,
            capture_output=True,
        )
        return png
    return path


def add_lead(prs, spec, footer: str, tmp: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_fill(slide, BRAND)
    geom = ROOT / "images" / "title-geometry.png"
    # Right-weighted cover so the left stays readable blue.
    add_picture_fit(slide, geom, Inches(4.6), Inches(0), height=SH)
    logo = ROOT / "images" / "dtd-logo-white.png"
    add_picture_fit(slide, logo, SW - Inches(2.35), Inches(0.22), width=Inches(1.88))

    thank = "thank-you-slide" in spec["classes"]
    title_size = 32 if thank else 30
    box, tf = add_textbox(slide, Inches(0.75), Inches(2.05 if thank else 1.85), Inches(8.6), Inches(1.6))
    p = clear_first_p(tf)
    run = p.add_run()
    run.text = spec["title"]
    set_run(run, size=title_size, color=WHITE, bold=True)

    y = Inches(3.55 if thank else 3.55)
    paras = paragraphs(spec["body"])
    if paras:
        box, tf = add_textbox(slide, Inches(0.75), y, Inches(8.4), Inches(1.5))
        for i, para in enumerate(paras):
            p = clear_first_p(tf) if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            add_inline(p, para, size=16 if thank else 16, color=WHITE)

    if thank:
        imgs = parse_images(spec["body"])
        qr_src = next((p for p, _ in imgs if "qr" in Path(p).name.lower()), None)
        pad = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.55), Inches(3.55), Inches(2.45)
        )
        pad.fill.solid()
        pad.fill.fore_color.rgb = WHITE
        pad.line.fill.background()
        try:
            pad.adjustments[0] = 0.08
        except Exception:
            pass
        if qr_src:
            add_picture_fit(
                slide,
                ROOT / qr_src,
                Inches(1.35),
                Inches(4.68),
                width=Inches(1.85),
                height=Inches(1.85),
            )
        url = "https://peppelinux.github.io/Wallet-Presentations/gdc-italian-federation-vc-wallet/"
        box, tf = add_textbox(slide, Inches(0.82), Inches(6.52), Inches(3.4), Inches(0.40))
        p = clear_first_p(tf)
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "peppelinux.github.io/Wallet-Presentations/…"
        set_run(run, size=8, color=BRAND, bold=True)
        run.hyperlink.address = url

    box, tf = add_textbox(slide, Inches(0.75), SH - Inches(0.42), Inches(10.5), Inches(0.28))
    p = clear_first_p(tf)
    run = p.add_run()
    run.text = footer
    set_run(run, size=9, color=RGBColor(0xC5, 0xDD, 0xF5))


def add_content_slide(prs, spec, footer: str, tmp: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    classes = spec["classes"]
    title_size = 22 if "lessons-slide" in classes or "italian-choice-slide" in classes else 24
    body_top = chrome(slide, spec["title"], footer, title_size=title_size)
    body = spec["body"]
    body_h = SH - body_top - Inches(0.48)

    if "italian-choice-slide" in classes:
        items = parse_lists(body)
        box, tf = add_textbox(slide, ML, body_top, CONTENT_W, body_h)
        add_bullets(tf, items, size=14, space_after=8)
        return

    if "federation-api-compact" in classes or "domestic-gap-matrix" in classes:
        y = body_top
        intro = paragraphs(body)
        if "domestic-gap-matrix" in classes and intro:
            box, tf = add_textbox(slide, ML, y, CONTENT_W, Inches(0.38))
            add_inline(clear_first_p(tf), intro[0], size=13, color=MUTED)
            y += Inches(0.40)
            intro = intro[1:]
        parsed = parse_table(body)
        if parsed:
            headers, rows = parsed
            n = 1 + len(rows)
            tbl_h = Inches(0.36 * n + 0.08)
            shape = slide.shapes.add_table(n, len(headers), ML, y, CONTENT_W, tbl_h)
            table = shape.table
            col0 = 0.24 if "domestic-gap-matrix" in classes else 0.28
            table.columns[0].width = int(CONTENT_W * col0)
            table.columns[1].width = int(CONTENT_W * (1 - col0))
            for i, h in enumerate(headers):
                fill_cell(table.cell(0, i), h, size=12, header=True)
            for r, row in enumerate(rows, start=1):
                for c, val in enumerate(row):
                    fill_cell(table.cell(r, c), val, size=12, bold=(c == 0))
            style_table(table, font=12)
            y += tbl_h + Inches(0.14)
        if intro:
            box, tf = add_textbox(slide, ML, y, CONTENT_W, SH - y - Inches(0.44))
            for i, para in enumerate(intro):
                p = clear_first_p(tf) if i == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                add_inline(p, para, size=13)
        return

    if "case-stats-slide" in classes:
        stats = parse_stats(body)
        gap = Inches(0.16)
        n = max(len(stats), 1)
        card_w = int((CONTENT_W - gap * (n - 1)) / n)
        x = ML
        card_h = Inches(1.28)
        for num, lab in stats:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, body_top, card_w, Emu(63500))
            bar.fill.solid()
            bar.fill.fore_color.rgb = BRAND
            bar.line.fill.background()
            card = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, body_top + Emu(63500), card_w, card_h - Emu(63500)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = BRAND_SOFT
            card.line.fill.background()
            box, tf = add_textbox(slide, x, body_top + Inches(0.18), card_w, Inches(0.55))
            p = clear_first_p(tf)
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = num
            set_run(run, size=26, color=BRAND, bold=True)
            box, tf = add_textbox(slide, x + Inches(0.08), body_top + Inches(0.72), card_w - Inches(0.16), Inches(0.48))
            p = clear_first_p(tf)
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = lab
            set_run(run, size=11, color=MUTED, bold=True)
            x += card_w + gap
        y = body_top + card_h + Inches(0.18)
        items = parse_lists(body)
        box, tf = add_textbox(slide, ML, y, CONTENT_W, Inches(1.55))
        add_bullets(tf, items, size=14, space_after=5)
        rest = paragraphs(body)
        if rest:
            box, tf = add_textbox(slide, ML, y + Inches(1.55), CONTENT_W, Inches(0.55))
            add_inline(clear_first_p(tf), rest[0], size=13)
        return

    if "lessons-slide" in classes:
        y = body_top
        sections = re.split(r"\n\*\*(What worked|What is hard)\*\*\n", body)
        # split keeps delimiters when using capturing groups
        parts = []
        buf = body
        for label in ("What worked", "What is hard"):
            m = re.search(rf"\*\*{re.escape(label)}\*\*", buf)
            if m:
                parts.append(label)
        blocks = re.split(r"\n\*\*(?:What worked|What is hard)\*\*\s*\n", body)
        # blocks[0] empty/preamble, then items for each section
        labeled = re.findall(
            r"\*\*(What worked|What is hard)\*\*\s*\n(.*?)(?=\n\*\*What |\Z)",
            body,
            re.S,
        )
        for heading, chunk in labeled:
            box, tf = add_textbox(slide, ML, y, CONTENT_W, Inches(0.32))
            p = clear_first_p(tf)
            run = p.add_run()
            run.text = heading
            set_run(run, size=15, color=BRAND, bold=True)
            y += Inches(0.30)
            items = parse_lists(chunk)
            h = Inches(0.36 * len(items) + 0.08)
            box, tf = add_textbox(slide, ML, y, CONTENT_W, h)
            add_bullets(tf, items, size=13, space_after=4)
            y += h + Inches(0.06)
        return

    if "roadmap-slide" in classes:
        items = parse_lists(body)
        box, tf = add_textbox(slide, ML, body_top, CONTENT_W, Inches(4.6))
        add_bullets(tf, items, size=15, space_after=8)
        rest = paragraphs(body)
        if rest:
            box, tf = add_textbox(slide, ML, Inches(6.35), CONTENT_W, Inches(0.55))
            add_inline(clear_first_p(tf), rest[0], size=13, color=MUTED)
        return

    # topology / runtime / browser / agenda (generic)
    imgs = parse_images(body)
    items = parse_lists(body)
    paras = paragraphs(body)
    numbered = bool(re.search(r"^\d+\. ", body, re.M))

    y = body_top
    if paras and (
        "browser-shot-slide" in classes
        or (not items and imgs)
        or "topology-slide" in classes
    ):
        # first paragraph sits above the figure on browser / topology caption after
        if "browser-shot-slide" in classes and paras:
            box, tf = add_textbox(slide, ML, y, CONTENT_W, Inches(0.42))
            add_inline(clear_first_p(tf), paras[0], size=13)
            y += Inches(0.40)
            paras = paras[1:]

    if imgs:
        src, href = imgs[0]
        img_path = resolve_asset(src, tmp)
        max_h = Inches(3.15 if items else 4.35)
        if "runtime-slide" in classes:
            max_h = Inches(2.85)
        if "browser-shot-slide" in classes:
            max_h = Inches(5.05)
        # keep aspect
        from PIL import Image

        with Image.open(img_path) as im:
            iw, ih = im.size
        aspect = iw / ih
        w = CONTENT_W
        h = int(w / aspect)
        if h > max_h:
            h = max_h
            w = int(h * aspect)
        left = ML + (CONTENT_W - w) // 2
        pic = add_picture_fit(slide, img_path, left, y, width=w, height=h)
        if href:
            pic.click_action.hyperlink.address = href
        y += h + Inches(0.10)

    if items:
        box, tf = add_textbox(slide, ML, y, CONTENT_W, SH - y - Inches(0.44))
        if numbered:
            for i, item in enumerate(items):
                p = clear_first_p(tf) if i == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                p.line_spacing = 1.2
                add_inline(p, f"{i + 1}. {item}", size=18 if len(items) <= 4 else 15)
        else:
            add_bullets(tf, items, size=14 if "runtime-slide" in classes else 15, space_after=6)
        y += Inches(0.2)

    if paras:
        # topology caption under diagram when bullets absent
        box, tf = add_textbox(slide, ML, y if not items else Inches(6.35), CONTENT_W, Inches(0.70))
        for i, para in enumerate(paras):
            p = clear_first_p(tf) if i == 0 else tf.add_paragraph()
            p.space_after = Pt(4)
            add_inline(p, para, size=13)


def build():
    meta, body = parse_front_matter(DECK.read_text(encoding="utf-8"))
    footer = meta.get(
        "footer",
        "GDC 2026 · Geneva · 2 September · OpenID Federation Part 1 · Case study · Giuseppe De Marco",
    )
    specs = split_slides(body)

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    # unused default layouts remain; we only use blank

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        for spec in specs:
            if "lead" in spec["classes"]:
                add_lead(prs, spec, footer, tmp_path)
            else:
                add_content_slide(prs, spec, footer, tmp_path)

    # drop unused unused first slide if any — we only added blanks
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides)} native slides)")


if __name__ == "__main__":
    build()
